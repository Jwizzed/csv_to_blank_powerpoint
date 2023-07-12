from pptx import Presentation
import pandas as pd


class BlankTextSlide:
    def __init__(self, csv_file):
        self.prs = Presentation()
        self.csv = pd.read_csv(csv_file, delimiter=",", na_values="nan")
        self.title_slide_layout = self.prs.slide_layouts[0]
        self.df = None

    def print_data_info(self):
        """Prints the data in the csv file"""
        print("Data in the csv file: \n")
        self.df = pd.DataFrame(self.csv).applymap(str)
        print(self.df)

    def create_slide(self):
        """Creates a blank slide for each row in the csv file"""
        for name, desc in self.df.values:
            slide = self.prs.slides.add_slide(self.title_slide_layout)
            title = slide.shapes.title
            title.text = name
            if desc != 'nan':
                slide.placeholders[1].text = desc
        self.prs.save('test.pptx')
        print("\nSlide created successfully")


if __name__ == '__main__':
    file_path = input("Enter the path to the csv file: ") # Ex: "example.csv"
    b = BlankTextSlide(file_path)
    b.print_data_info()
    b.create_slide()